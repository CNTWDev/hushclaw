"""PPT editing tools powered by python-pptx.

All functions are synchronous (python-pptx is a sync library).
Installed automatically when the skill is loaded via HushClaw's bundled-tool mechanism.
"""
from __future__ import annotations
import json

from pathlib import Path

from hushclaw.tools.base import ToolResult, tool


def _split_lines(text: str) -> list[str]:
    lines = [x.strip() for x in text.splitlines() if x.strip()]
    return lines


def _as_title_case(text: str) -> str:
    return str(text or "").strip() or "Untitled"


@tool(description="Get slide count and title list from a PPTX file.")
def pptx_info(path: str) -> ToolResult:
    """Return basic info about the presentation: slide count and per-slide titles."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ToolResult(error="python-pptx is not installed. Run: pip install python-pptx")

    p = Path(path).expanduser()
    if not p.exists():
        return ToolResult(error=f"File not found: {path}")

    prs = Presentation(str(p))
    titles = []
    for i, slide in enumerate(prs.slides):
        title_shape = slide.shapes.title
        titles.append({
            "slide": i,
            "title": title_shape.text.strip() if title_shape and title_shape.has_text_frame else "",
        })

    return ToolResult.ok(json.dumps({
        "path": str(p),
        "slide_count": len(prs.slides),
        "slides": titles,
    }, ensure_ascii=False))


@tool(description="Read all text from a specific slide (0-based index).")
def pptx_read_slide(path: str, slide_index: int) -> ToolResult:
    """Extract all text from the specified slide."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ToolResult(error="python-pptx is not installed. Run: pip install python-pptx")

    p = Path(path).expanduser()
    if not p.exists():
        return ToolResult(error=f"File not found: {path}")

    prs = Presentation(str(p))
    if slide_index < 0 or slide_index >= len(prs.slides):
        return ToolResult(error=f"slide_index {slide_index} out of range (0–{len(prs.slides)-1})")

    slide = prs.slides[slide_index]
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    texts.append(line)

    return ToolResult.ok(json.dumps({
        "slide": slide_index,
        "texts": texts,
        "full_text": "\n".join(texts),
    }, ensure_ascii=False))


@tool(description="Extract all text from every slide in a PPTX file.")
def pptx_extract_all_text(path: str) -> ToolResult:
    """Return a list of {slide, texts} dicts covering all slides."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ToolResult(error="python-pptx is not installed. Run: pip install python-pptx")

    p = Path(path).expanduser()
    if not p.exists():
        return ToolResult(error=f"File not found: {path}")

    prs = Presentation(str(p))
    slides_out = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        slides_out.append({"slide": i, "texts": texts})

    return ToolResult.ok(json.dumps({"path": str(p), "slides": slides_out}, ensure_ascii=False))


@tool(description="Add a title slide (with title and subtitle) to an existing PPTX file.")
def pptx_add_title_slide(path: str, title: str, subtitle: str = "") -> ToolResult:
    """Append a title slide layout page to the presentation."""
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Pt  # type: ignore  # noqa: F401
    except ImportError:
        return ToolResult(error="python-pptx is not installed. Run: pip install python-pptx")

    p = Path(path).expanduser()
    if not p.exists():
        return ToolResult(error=f"File not found: {path}")

    prs = Presentation(str(p))
    layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(layout)

    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    if 0 in placeholders:
        placeholders[0].text = title
    if 1 in placeholders:
        placeholders[1].text = subtitle

    prs.save(str(p))
    return ToolResult.ok(json.dumps({
        "added_slide": len(prs.slides) - 1,
        "title": title,
        "subtitle": subtitle,
    }, ensure_ascii=False))


@tool(description="Add a text content slide (title + body text) to an existing PPTX file.")
def pptx_add_text_slide(path: str, title: str, content: str) -> ToolResult:
    """Append a title+content layout slide. Use newlines in content for bullet points."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ToolResult(error="python-pptx is not installed. Run: pip install python-pptx")

    p = Path(path).expanduser()
    if not p.exists():
        return ToolResult(error=f"File not found: {path}")

    prs = Presentation(str(p))
    layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(layout)

    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    if 0 in placeholders:
        placeholders[0].text = title
    if 1 in placeholders:
        tf = placeholders[1].text_frame
        tf.clear()
        lines = content.splitlines()
        for i, line in enumerate(lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                tf.add_paragraph().text = line

    prs.save(str(p))
    return ToolResult.ok(json.dumps({
        "added_slide": len(prs.slides) - 1,
        "title": title,
    }, ensure_ascii=False))


@tool(description="Set the text of a specific placeholder on a slide (0-based slide_index and placeholder_index).")
def pptx_set_slide_text(path: str, slide_index: int, placeholder_index: int, text: str) -> ToolResult:
    """Modify the text of one placeholder on the specified slide."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ToolResult(error="python-pptx is not installed. Run: pip install python-pptx")

    p = Path(path).expanduser()
    if not p.exists():
        return ToolResult(error=f"File not found: {path}")

    prs = Presentation(str(p))
    if slide_index < 0 or slide_index >= len(prs.slides):
        return ToolResult(error=f"slide_index {slide_index} out of range (0–{len(prs.slides)-1})")

    slide = prs.slides[slide_index]
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    if placeholder_index not in placeholders:
        available = sorted(placeholders.keys())
        return ToolResult(error=f"placeholder_index {placeholder_index} not found. Available: {available}")

    placeholders[placeholder_index].text = text
    prs.save(str(p))
    return ToolResult.ok(json.dumps({
        "slide": slide_index,
        "placeholder": placeholder_index,
        "text": text,
    }, ensure_ascii=False))


@tool(description="Delete a slide from a PPTX file by 0-based slide_index.")
def pptx_delete_slide(path: str, slide_index: int) -> ToolResult:
    """Remove the slide at slide_index from the presentation."""
    try:
        from pptx import Presentation  # type: ignore
        from lxml import etree  # type: ignore  # noqa: F401
    except ImportError as e:
        return ToolResult(error=f"Missing dependency: {e}. Run: pip install python-pptx lxml")

    p = Path(path).expanduser()
    if not p.exists():
        return ToolResult(error=f"File not found: {path}")

    prs = Presentation(str(p))
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        return ToolResult(error=f"slide_index {slide_index} out of range (0–{total-1})")

    xml_slides = prs.slides._sldIdLst
    xml_slides.remove(xml_slides[slide_index])
    prs.save(str(p))
    return ToolResult.ok(json.dumps({
        "deleted_slide": slide_index,
        "remaining_slides": total - 1,
    }, ensure_ascii=False))


@tool(description="Create a new blank PPTX file at the given path. Overwrites if exists.")
def pptx_create(path: str) -> ToolResult:
    """Create an empty PowerPoint presentation."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ToolResult(error="python-pptx is not installed. Run: pip install python-pptx")

    p = Path(path).expanduser()
    p.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.save(str(p))
    return ToolResult.ok(json.dumps({"created": str(p)}, ensure_ascii=False))


@tool(
    description=(
        "Add a consulting-style insight slide with visual hierarchy: "
        "answer-first headline, icon cards, and so-what action box."
    )
)
def pptx_add_consulting_insight_slide(
    path: str,
    headline: str,
    key_points: str,
    so_what: str = "",
    icon: str = "◆",
) -> ToolResult:
    """Append a designed slide that is visually richer than the basic text layout."""
    try:
        from pptx import Presentation  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
    except ImportError:
        return ToolResult(error="python-pptx is not installed. Run: pip install python-pptx")

    p = Path(path).expanduser()
    if not p.exists():
        return ToolResult(error=f"File not found: {path}")

    points = _split_lines(key_points)
    if len(points) == 0:
        return ToolResult(error="key_points must contain at least one non-empty line")
    points = points[:4]

    prs = Presentation(str(p))
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Accent bar
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.4),
        Inches(0.45),
        Inches(0.08),
        Inches(6.5),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(39, 94, 254)
    accent.line.fill.background()

    # Headline
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(1.0))
    tf = title_box.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = headline
    p0.font.bold = True
    p0.font.size = Pt(30)
    p0.font.color.rgb = RGBColor(30, 41, 59)

    # Cards
    card_y = 1.8
    card_h = 1.05
    gap = 0.25
    card_w = (11.8 - gap * (len(points) - 1)) / len(points)
    start_x = 0.7

    for idx, point in enumerate(points):
        x = start_x + idx * (card_w + gap)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(card_y),
            Inches(card_w),
            Inches(card_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 248, 255)
        card.line.color.rgb = RGBColor(214, 224, 255)
        ctf = card.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.text = f"{icon}  {point}"
        cp.font.size = Pt(16)
        cp.font.bold = True
        cp.font.color.rgb = RGBColor(37, 52, 84)

    # So-what/action box
    action_text = so_what.strip() or "So what: insert decision and immediate action."
    action = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.7),
        Inches(3.25),
        Inches(11.8),
        Inches(2.9),
    )
    action.fill.solid()
    action.fill.fore_color.rgb = RGBColor(255, 255, 255)
    action.line.color.rgb = RGBColor(226, 232, 240)
    atf = action.text_frame
    atf.clear()
    ap = atf.paragraphs[0]
    ap.text = "So what / Decision"
    ap.font.bold = True
    ap.font.size = Pt(18)
    ap.font.color.rgb = RGBColor(30, 41, 59)
    ap2 = atf.add_paragraph()
    ap2.text = action_text
    ap2.font.size = Pt(15)
    ap2.font.color.rgb = RGBColor(51, 65, 85)

    prs.save(str(p))
    return ToolResult(
        output={
            "added_slide": len(prs.slides) - 1,
            "headline": headline,
            "point_count": len(points),
            "icon": icon,
        }
    )


@tool(
    description=(
        "Add a premium consulting template slide. template supports: "
        "strategy_house, matrix_2x2, waterfall, timeline."
    )
)
def pptx_add_consulting_template_slide(
    path: str,
    template: str,
    title: str,
    content_lines: str,
    icon: str = "◆",
) -> ToolResult:
    """Append one of several high-fidelity consulting slide templates."""
    try:
        from pptx import Presentation  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
    except ImportError:
        return ToolResult(error="python-pptx is not installed. Run: pip install python-pptx")

    p = Path(path).expanduser()
    if not p.exists():
        return ToolResult(error=f"File not found: {path}")

    tpl = str(template or "").strip().lower()
    if tpl not in {"strategy_house", "matrix_2x2", "waterfall", "timeline"}:
        return ToolResult(error="template must be one of: strategy_house, matrix_2x2, waterfall, timeline")

    lines = _split_lines(content_lines)
    prs = Presentation(str(p))
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # shared title zone
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = _as_title_case(title)
    p0.font.bold = True
    p0.font.size = Pt(30)
    p0.font.color.rgb = RGBColor(30, 41, 59)

    if tpl == "strategy_house":
        # roof
        roof = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(3.5), Inches(1.2), Inches(5.4), Inches(1.2))
        roof.fill.solid()
        roof.fill.fore_color.rgb = RGBColor(39, 94, 254)
        roof.line.fill.background()
        rtf = roof.text_frame
        rtf.text = f"{icon}  {(lines[0] if len(lines) > 0 else 'Aspiration')}"
        rtf.paragraphs[0].font.size = Pt(16)
        rtf.paragraphs[0].font.bold = True
        rtf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        # pillars
        for i in range(3):
            x = 2.5 + i * 2.0
            pillar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(2.5), Inches(1.4), Inches(2.2))
            pillar.fill.solid()
            pillar.fill.fore_color.rgb = RGBColor(245, 248, 255)
            pillar.line.color.rgb = RGBColor(203, 213, 225)
            txt = lines[i + 1] if len(lines) > i + 1 else f"Initiative {i+1}"
            pillar.text_frame.text = txt
            pillar.text_frame.paragraphs[0].font.size = Pt(14)
            pillar.text_frame.paragraphs[0].font.bold = True
        # foundation
        base = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.0), Inches(5.0), Inches(7.6), Inches(0.6))
        base.fill.solid()
        base.fill.fore_color.rgb = RGBColor(226, 232, 240)
        base.line.fill.background()
        base.text_frame.text = lines[4] if len(lines) > 4 else "Foundation / Capability / Governance"
        base.text_frame.paragraphs[0].font.size = Pt(13)

    elif tpl == "matrix_2x2":
        # axes
        h = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.2), Inches(5.3), Inches(8.2), Inches(0.05))
        h.fill.solid()
        h.fill.fore_color.rgb = RGBColor(100, 116, 139)
        h.line.fill.background()
        v = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.2), Inches(1.4), Inches(0.05), Inches(4.0))
        v.fill.solid()
        v.fill.fore_color.rgb = RGBColor(100, 116, 139)
        v.line.fill.background()
        # quadrant labels
        labels = (lines + ["Q1", "Q2", "Q3", "Q4"])[:4]
        pos = [(2.5, 1.8), (6.6, 1.8), (2.5, 3.8), (6.6, 3.8)]
        for (x, y), lab in zip(pos, labels):
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3.2), Inches(1.2))
            box.text_frame.text = f"{icon}  {lab}"
            box.text_frame.paragraphs[0].font.bold = True
            box.text_frame.paragraphs[0].font.size = Pt(15)

    elif tpl == "waterfall":
        vals = (lines + ["Step 1", "Step 2", "Step 3", "Result"])[:4]
        heights = [1.3, 0.9, 1.6, 2.2]
        colors = [RGBColor(148, 163, 184), RGBColor(99, 102, 241), RGBColor(37, 99, 235), RGBColor(16, 185, 129)]
        x = 1.4
        for i, (label, hgt) in enumerate(zip(vals, heights)):
            bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(5.4 - hgt), Inches(2.1), Inches(hgt))
            bar.fill.solid()
            bar.fill.fore_color.rgb = colors[i]
            bar.line.fill.background()
            bar.text_frame.text = label
            bar.text_frame.paragraphs[0].font.size = Pt(13)
            bar.text_frame.paragraphs[0].font.bold = True
            bar.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            x += 2.5

    elif tpl == "timeline":
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.2), Inches(3.2), Inches(10.5), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(100, 116, 139)
        line.line.fill.background()
        marks = (lines + ["Phase 1", "Phase 2", "Phase 3", "Phase 4"])[:4]
        for i, m in enumerate(marks):
            x = 1.4 + i * 2.6
            dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(3.0), Inches(0.35), Inches(0.35))
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(37, 99, 235)
            dot.line.fill.background()
            lab = slide.shapes.add_textbox(Inches(x - 0.1), Inches(3.45), Inches(2.0), Inches(1.0))
            lab.text_frame.text = f"{icon}  {m}"
            lab.text_frame.paragraphs[0].font.size = Pt(13)
            lab.text_frame.paragraphs[0].font.bold = True

    prs.save(str(p))
    return ToolResult.ok(json.dumps({"added_slide": len(prs.slides) - 1, "template": tpl, "title": title}, ensure_ascii=False))
