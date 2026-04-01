"""Icon embedding tools for HushClaw PPT skill.

Bundles 61 Phosphor Icons (MIT) SVGs from assets/icons/.
Converts SVG → PNG in memory using cairosvg (if installed) or
falls back to a Pillow-drawn geometric shape that approximates
the icon's visual category.

No external network calls are made at runtime — all icons are local.
"""
from __future__ import annotations

import base64
import io
import json
import os
import tempfile
from pathlib import Path
from typing import Optional

from hushclaw.tools.base import ToolResult, tool

_ICONS_DIR = Path(__file__).resolve().parent.parent / "assets" / "icons"
_INDEX_FILE = _ICONS_DIR / "_index.json"

# ── Placeholder (gray gradient, 640×360, ~5 KB JPEG, no network needed) ──────
_PLACEHOLDER_B64 = (
    "/9j/4AAQSkZJRgABAQAAAQABAAD/2wBDAAkGBwgHBgkIBwgKCgkLDRYPDQwMDRsUFRAWIB0iIiAdHx8kKDQsJCYx"
    "Jx8fLT0tMTU3Ojo6Iys/RD84QzQ5Ojf/2wBDAQoKCg0MDRoPDxo3JR8lNzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3"
    "Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzf/wAARCAFoAoADASIAAhEBAxEB/8QAGwABAQEBAQEBAQAAAAAAAAAA"
    "AAIBBAMGBQf/xAA9EAEAAQMEAgECAgYIAwkAAAAAcjKx4QECEbIDMwQhMRJBBRMUIlFhFVJTcZGSodEjgfA0QkVk"
    "goOiwcL/xAAXAQEBAQEAAAAAAAAAAAAAAAAAAQID/8QAGBEBAQEBAQAAAAAAAAAAAAAAABEBAhL/2gAMAwEAAhED"
    "EQA/AP61TGyhNMbOjCgAE0xsoAE0xsoAAE0xsoTTGwKAATTGygATTGygAATTGyhNMbAoABNMbKABNMbKAABNMbKE"
    "0xsCgAE0xsoAE0xsoAAE0xsoTTGwKAATTGygATTGygAATTGyhNMbAoABNMbKABNMbKAABNMbKE0xsCgAE0xsoAE0"
    "xsoAAE0xsoTTGwKAATTGygATTGygAAA4OATTGyjhNMbAoODgBNMbK4OABNMbK4ADg4BNMbKOE0xsCg4OAE0xsrg4"
    "AE0xsrgAODgE0xso4TTGwKDg4ATTGyuDgATTGyuAA4OATTGyjhNMbAoODgBNMbK4OABNMbK4ADg4BNMbKOE0xsCg"
    "4OAE0xsrg4AE0xsrgAODgE0xso4TTGwKDg4ATTGyuDgATTGyuAA4OATTGyjhNMbAoODgBNMbK4OABNMbK4ADg4BN"
    "MbKOE0xsCg4OAE0xsrg4AE0xsrgAZTGzQAAZTGzRlMbA0ABlMbNABlMbNAABlMbNGUxsDQAGUxs0AGUxs0AAGUxs"
    "0ZTGwNAAZTGzQAZTGzQAAZTGzRlMbA0ABlMbNABlMbNAABlMbNGUxsDQAGUxs0AGUxs0AAGUxs0ZTGwNAAZTGzQA"
    "ZTGzQAAZTGzRlMbA0ABlMbNABlMbNAABlMbNGUxsDQAGUxs0AZTGygGBTGzQYNATTGzWspjYAaAxlMbKAYFMbNBg"
    "0BNMbNaymNgBoDGUxsoBgUxs0GDQE0xs1rKY2AGgMZTGygGBTGzQYNATTGzWspjYAaAxlMbKAYFMbNBg0BNMbNay"
    "mNgBoDGUxsoBgUxs0GDQE0xs1rKY2AGgMZTGygGBTGzQYNATTGzWspjYAaAxlMbKAYFMbNBg0BNMbNaymNgBoAAA"
    "ymNmgAymNmgAAymNmjKY2BoADKY2aADKY2aAADKY2aMpjYGgAMpjZoAMpjZoAAMpjZoymNgaAAymNmgAymNmgAAy"
    "mNmjKY2BoADKY2aADKY2aAADKY2aMpjYGgAMpjZoAMpjZoAAMpjZoymNgaAAymNmgAymNmgAAymNmjKY2BoADKY2"
    "aADKY2aAADKY2aMpjYCmNlCaY2BQACaY2UACaY2UAACaY2UJpjYFAAJpjZQAJpjZQAAJpjZQmmNgUAAmmNlAAmmN"
    "lAAAmmNlCaY2BQACaY2UACaY2UAACaY2UJpjYFAAJpjZQAJpjZQAAJpjZQmmNgUAAmmNlAAmmNlAAAmmNlCaY2BQ"
    "ACaY2UACaY2UAACaY2UJpjYFAAJpjZQAJpjZQAAACKmmNlCaY2EUAKJpjZQAJpjZQAAJpjZQmmNhFACiaY2UACaY"
    "2UAACaY2UJpjYRQAommNlAAmmNlAAAmmNlCaY2EUAKJpjZQAJpjZQAAJpjZQmmNhFACiaY2UACaY2UAACaY2UJpj"
    "YRQAommNlAAmmNlAAAmmNlCaY2EUAKJpjZQAJpjZQAAJpjZQmmNhFACiaY2UACaY2UDRNEbKAABNEbKE0RsIoAUT"
    "RGygATRGygAATRGyhNEbCKAFE0RsoAE0RsoAAE0RsoTRGwigBRNEbKABNEbKAABNEbKE0RsIoAUTRGygATRGygAA"
    "TRGyhNEbCKAFE0RsoAE0RsoAAE0RsoTRGwigBRNEbKABNEbKAABNEbKE0RsIoAUTRGygATRGygAATRGyhNEbCKAF"
    "E0RsoATRGywGDKI2UDBoCKI2U1NEbCNGgrE0RssBgyiNlAwaAiiNlNTRGwjRoKxNEbLAYMojZQMGgIojZTU0RsI0"
    "aCsTRGywGDKI2UDBoCKI2U1NEbCNGgrE0RssBgyiNlAwaAiiNlNTRGwjRoKxNEbLAYMojZQMGgIojZTU0RsI0aCs"
    "TRGywGDKI2UDBoCKI2U1NEbCNGgrE0RssBgyiNlAwaAiiNlNTRGwjRoKAAJojbCgATRG2FAAAmiNsKE0RtgRQAom"
    "iNsKABNEbYUAACaI2woTRG2BFACiaI2woAE0RthQAAJojbChNEbYEUAKJojbCgATRG2FAPD5Hydvh3bdmm3Xf5d9"
    "Ozb9/wDn/DRnyvN5Nuu3w/H00182/wDPX7bdP46q+N8fb8fZrpprru3bvru36/fdqqPD9V8rya/8Xy/qvFr/ANzx"
    "ffT/AJr/AKP8G76+X9Z5d39bfv111/udSaI2wUc/9H/H0+vj27vHu/Lds3686M18HyfFrp+z/I136fnt831/1+7r"
    "Cjn8PytN/k/U+XZr4vLxz+HXXnTX+7X83tRG2E+fw7PP49fH5NOdNf8AHT+bx+N5PL4/J+zfI153cc7PJ/X0/wBw"
    "dQmiNsKRQAE0RthQmiNsCKAFE0RthQAJojbCgAATRG2FCaI2wIoAUTRG2FAAmiNsKAABNEbYUJojbAigBRNEbYUA"
    "CaI2woAAE0RthQmiNsCFEbYU1NEbYBo0FYmiNsLAYMojbCgYNARRG2FNTRG2BGjQViaI2wsBgyiNsKBg0BFEbYU1"
    "NEbYEaNBWJojbCwGDKI2woGDQEURthTU0RtgRo0FY8/Lv08Hj3b91G3Tn+56uT9IfvbfD4vvp5PLppu2/nrt++pi"
    "N+D49dPHr5t/18nm/e3a8/Tj8tP8HSyiNsKKMfgfP+Zv+R5N23bu108Wn0000/P+er6B858z42743m1266a/h1o1"
    "/jo1yaz4vyvJ8bfprs1112/ns11+mr6LZu037Nu/brzt3ac6a/yfNeDw7/keTTx+PTnXX/DTT+L6Px7dPDs27NNd"
    "ddm3TTTnX76HSYt4/K8H6/w67dNddN2n72zXTXjjd+T3Ga08PiefT5Hx9vk0451+m7TT8tXpRG2HP8b9z5nyfH9t"
    "uuum/bpr+fOn11/xdZqMGURthQrBoCKI2wpqaI2wI0aCsTRG2FgMGURthQMGgIojbCmpojbAjRoKxNEbYWAwZRG2"
    "FAwaAiiNsKamiNsCNGgrE0RthYDBlEbYUDBoAACaI2woTRG2BFACiaI2woAE0RthQAAJojbChNEbYEUAKJojbCgA"
    "TRG2FAAAmiNsKE0RtgRQAomiNsKABNEbYUAACaI2woTRG2BFOT5P/b/h/wDrs63J8/8Ac3fH8un012eXTTXd+Wm3"
    "X6arg600RthQijN+3bv267d+3Tdt1++munOjKI2woE7PHs8enHj2bdun3426cKAE0RthQmiNsCOb/wAV/wDY/wD0"
    "63J4P3/nfI36/XTZpt2bdfy0/PXT/F1roJojbChFBNEbYUAACaI2woTRG2BFACiaI2woAE0RthQAAJojbChNEbYE"
    "UAKJojbCgATRG2FAAAmiNsKE0RtgRQAomiNsKABNEbYUDRFEbYWAACKI2wsRRG2AWAAiiNsLABFEbYWAACKI2wsR"
    "RG2AWAAiiNsLABFEbYWAACKI2wsRRG2AWAAiiNsLABFEbYWAACKI2wz5Hh2+fw7/ABbvtu0+/wDDV6IojbAPH4Pl"
    "/WeHTbu1/wCL4/3d+muvOumujpcvytnk8e/T5Hxtumu7T2bP6+n++n/X8HJ+kP0jt3fH27Pj66a/rNNfxc/fbp/B"
    "ZUqvL+l9uzy7tuzxfj26a8abvx/f/R57f0xx9NPB9Py0/H9v9H5Q1MSv1v6Z/wDL/wDzwafpnTnT8Xg100/PXTfz"
    "/wDT8kJhdfV6a6btNNduumumv1010/N5/J82z4/h3eTfrp9Ptpz99f4Py/0X87TxbdfD5t2mmzTTXXbu/h/J2+DT"
    "f8ry6efy7fw+Hb9fFs1+/P8AW1/6zmRar4Xh1+N4NNN/113fvb9f4a6uoRRG2EFgCiKI2wsAEURthYAAIojbCxFE"
    "bYBYACKI2wsAEURthYAAIojbCxFEbYBYACKI2wsAEURthYAAIojbCxFEbYBYACKI2wsARRG2FiAJojbCgAARRG2F"
    "iaI2wooBARRG2FgAmiNsKAABFEbYWJojbCigEBFEbYWACaI2woAAEURthYmiNsKKAQEURthYAJojbCgAAc/yvPt+"
    "J4vx7tNddv227dP4/wCz53fu13792/drzu3a866/zfT+TxePycfrPHs38fb8WnPDz1+L4NuvOvg8X4YafRrNibj5"
    "ofT/ALL8f+w8X+TQ/Zfj/wBh4v8AJovpI+YH0/7L8f8AsPF/k0P2X4/9h4v8mh6I+YfRfo/5enyvD9ef1mzjTf8A"
    "z/mvX4vg2686+Dxfhhp9Hps8Hi8evPj8ezbr9udu3TRN2rmLAZVFEbYWJojbCigEBFEbYWACaI2woAAEURthYmiN"
    "sKKAQEURthYAJojbCgAARRG2FiaI2wooBARRG2FgAmiNsKAABFEbYWJojbCigEGgAIojbCwARRG2FgAAiiNsLEUR"
    "thUWAiiKI2wsAEURthYAAIojbCxFEbYVFgIoiiNsLABFEbYWAACKI2wsRRG2FRYCKIojbCwARRG2FgAAiiNsLEUR"
    "thUWAiiKI2wsAEURthYAAIojbCxFEbYVFgIoiiNsLABFEbYWAACKI2wsRRG2FRYCKIojbCwARRG2FgAAiiNsLEUR"
    "thUWAiiKI2wsAEURthYAAIojbCxFEbYVCiNsLE0QtgVQCAiiNsLABNELYUAACKI2wsTRC2FFAICKI2wsAE0QthQA"
    "AIojbCxNELYUUAgIojbCwATRC2FAAAiiNsLE0QthRQCAiiNsLABNELYUAACKI2wsTRC2FFAICKI2wsAE0QthQAAI"
    "ojbCxNELYUUAgIojbCwATRC2FAAAiiNsLE0QthRQCAiiNsLABNELYUAACKI2wsTRC2FFAICKI2wsAE0QthQAANAB"
    "FELYWIohbCiwEBFELYWACKIWwsAAEUQthYiiFsKLAQEUQthYAIohbCwAARRC2FiKIWwosBARRC2FgAiiFsLAABFE"
    "LYWIohbCiwEBFELYWACKIWwsAAEUQthYiiFsKLAQEUQthYAIohbCwAARRC2FiKIWwosBARRC2FgAiiFsLAABFELY"
    "WIohbCiwEBFELYWACKIWwsAAEUQthYiiFsKLAQEUQthYAIohbCwBNEOuFAAAiiFsLE0Q64UUAgIohbCwATRDrhQA"
    "AIohbCxNEOuFFAICKIWwsAE0Q64UAACKIWwsTRDrhRQCAiiFsLABNEOuFAAAiiFsLE0Q64UUAgIohbCwATRDrhQA"
    "AIohbCxNEOuFFAICKIWwsAE0Q64UAACKIWwsTRDrhRQCAiiFsLABNEOuFAAAiiFsLE0Q64UUAgIohbCwATRDrhQA"
    "AIohbCxNEOuFFAICKIWwsATRDrhQAJoh1woAAE0Q64UJoh1wooBATRDrhQAJoh1woAAE0Q64UJoh1wooBATRDrhQ"
    "AJoh1woAAE0Q64UJoh1wooBATRDrhQAJoh1woAAE0Q64UJoh1wooBATRDrhQAJoh1woAAE0Q64UJoh1wooBATRDr"
    "hQAJoh1woAAE0Q64UJoh1wooBATRDrhQAJoh1woAAE0Q64UJoh1wooBATRDrhQAJoh1woAAE0Q64UJoh1wooBBoA"
    "CfXDrhQAJ9cOuFAAAn1w64UJ9cOuAUAAn1w64UACfXDrhQAAJ9cOuFCfXDrgFAAJ9cOuFAAn1w64UAACfXDrhQn1"
    "w64BQACfXDrhQAJ9cOuFAAAn1w64UJ9cOuAUAAn1w64UACfXDrhQAAJ9cOuFCfXDrgFAAJ9cOuFAAn1w64UAACfX"
    "DrhQn1w64BQACfXDrhQAJ9cOuFAAAn1w64UJ9cOuAUAAn1w64UACfXDrhQAAJ9cOuFCfXDrgD1w64U3hHrh1wChv"
    "BwDE+uHXC+DgGCfXDrhfAMG8HAI9cOuFN4R64dcAobwcAxPrh1wvg4Bgn1w64XwDBvBwCPXDrhTeEeuHXAKG8HAM"
    "T64dcL4OAYJ9cOuF8AwbwcAj1w64U3hHrh1wChvBwDE+uHXC+DgGCfXDrhfAMG8HAI9cOuFN4R64dcAobwcAxPrh"
    "1wvg4Bgn1w64XwDBvBwCPXDrhTeEeuHXAKG8HAMT64dcL4OAYJ9cOuF8AwbwcAj1w64U3hHrh1wChvBwDE+uHXC+"
    "DgGCfXDrhfAMG8HAI9cOuFN4R64dcAobwcAxPrh1wvg4Bgn1w64XwDBvBwAACPXDrhYj1w64BYACPXDrhYAI9cOu"
    "FgAAj1w64WI9cOuAWAAj1w64WACPXDrhYAAI9cOuFiPXDrgFgAI9cOuFgAj1w64WAACPXDrhYj1w64BYACPXDrhY"
    "AI9cOuFgAAj1w64WI9cOuAWAAj1w64WACPXDrhYAAI9cOuFiPXDrgFgAI9cOuFgAj1w64WAACPXDrhYj1w64BYAC"
    "PXDrhYAI9cOuFgAAj1w64WI9cOuAWAAj1w64WACPXDrhYAAAAI9cOuFgAAAj1w64AFgAAAj1w64WAAACPXDrgAWA"
    "AACPXDrhYAAAI9cOuABYAAAI9cOuFgAAAj1w64AFgAAAj1w64WAAACPXDrgAWAAACPXDrhYAAAI9cOuABYAAAI9c"
    "OuFgAAAj1w64AFgAAAj1w64WAAACPXDrgAf/2Q=="
)


# ── Icon index (name → filename) ──────────────────────────────────────────────

def _load_index() -> dict[str, str]:
    if _INDEX_FILE.exists():
        return json.loads(_INDEX_FILE.read_text(encoding="utf-8"))
    return {}


# ── SVG → PNG conversion ──────────────────────────────────────────────────────

def _svg_to_png_bytes(svg_path: Path, size: int, color_hex: str) -> bytes | None:
    """Convert a local SVG to PNG bytes, coloring strokes/fills with color_hex.

    Tries cairosvg first; falls back to Pillow-drawn geometric shape.
    """
    # -- Attempt 1: cairosvg (best quality) --
    try:
        import cairosvg  # type: ignore  # noqa: F401
        svg_text = svg_path.read_text(encoding="utf-8")
        # Recolor: replace currentColor / inherit with the requested hex
        svg_text = svg_text.replace('color="currentColor"', f'color="{color_hex}"')
        svg_text = svg_text.replace("currentColor", color_hex)
        png_bytes = cairosvg.svg2png(
            bytestring=svg_text.encode(),
            output_width=size,
            output_height=size,
        )
        return png_bytes
    except Exception:
        pass

    # -- Attempt 2: Pillow (geometric fallback) --
    return _pillow_icon_png(svg_path.stem, size, color_hex)


def _pillow_icon_png(icon_stem: str, size: int, color_hex: str) -> bytes:
    """Draw a simple geometric icon using Pillow as a fallback."""
    try:
        from PIL import Image, ImageDraw  # type: ignore
    except ImportError:
        return b""

    # Parse color
    ch = color_hex.lstrip("#")
    try:
        r, g, b = int(ch[0:2], 16), int(ch[2:4], 16), int(ch[4:6], 16)
    except Exception:
        r, g, b = 29, 78, 216

    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    c = size // 2
    m = size // 8  # margin
    lw = max(2, size // 16)

    # Category-based geometric shapes
    name = icon_stem.lower()
    if any(k in name for k in ("chart", "bar", "stat", "trend")):
        # Bar chart: 3 bars of different heights
        bw = (size - m * 4) // 3
        for i, h in enumerate([0.5, 0.8, 0.6]):
            bh = int((size - m * 2) * h)
            bx = m + i * (bw + m)
            draw.rectangle([(bx, size - m - bh), (bx + bw, size - m)], fill=(r, g, b, 220))
    elif any(k in name for k in ("target", "cross", "bullseye")):
        # Bullseye
        for ri in [c - m, c - m * 2, c - m * 3]:
            draw.ellipse([(c - ri, c - ri), (c + ri, c + ri)], outline=(r, g, b, 220), width=lw)
        draw.ellipse([(c - m, c - m), (c + m, c + m)], fill=(r, g, b, 220))
    elif any(k in name for k in ("check", "seal", "trophy")):
        # Checkmark
        draw.line([(m, c), (c - m, size - m * 2)], fill=(r, g, b, 220), width=lw * 2)
        draw.line([(c - m, size - m * 2), (size - m, m)], fill=(r, g, b, 220), width=lw * 2)
    elif any(k in name for k in ("lightbulb", "idea", "brain", "sparkle")):
        # Lightbulb
        draw.ellipse([(m * 2, m), (size - m * 2, c + m)], outline=(r, g, b, 220), width=lw)
        draw.rectangle([(c - m, c + m), (c + m, size - m * 2)], outline=(r, g, b, 220), width=lw)
        draw.line([(c - m, size - m * 2), (c + m, size - m * 2)], fill=(r, g, b, 220), width=lw)
    elif any(k in name for k in ("user", "person", "people")):
        # Person silhouette
        draw.ellipse([(c - m * 2, m), (c + m * 2, m * 4)], fill=(r, g, b, 220))
        draw.arc([(m, c), (size - m, size - m)], 180, 0, fill=(r, g, b, 220), width=lw * 2)
    elif any(k in name for k in ("rocket", "lightning", "arrow")):
        # Arrow up-right
        pts = [(c, m), (size - m, m), (size - m, c), (c + m, c - m), (size - m, m), (c + m, c - m)]
        draw.line(pts[:3], fill=(r, g, b, 220), width=lw * 2)
        draw.line([(size - m, m), (m, size - m)], fill=(r, g, b, 220), width=lw * 2)
    elif any(k in name for k in ("globe", "world", "network")):
        # Globe
        draw.ellipse([(m, m), (size - m, size - m)], outline=(r, g, b, 220), width=lw)
        draw.line([(c, m), (c, size - m)], fill=(r, g, b, 220), width=lw)
        draw.arc([(m, m), (size - m, size - m)], 0, 180, fill=(r, g, b, 220), width=lw)
    elif any(k in name for k in ("shield", "lock", "secure")):
        # Shield
        pts = [(c, m), (size - m, m * 2), (size - m, c + m), (c, size - m), (m, c + m), (m, m * 2), (c, m)]
        draw.polygon(pts, outline=(r, g, b, 220), width=lw)
    elif any(k in name for k in ("gear", "cog", "setting", "wrench")):
        # Gear (simplified)
        draw.ellipse([(c - m * 2, c - m * 2), (c + m * 2, c + m * 2)], outline=(r, g, b, 220), width=lw)
        for angle in range(0, 360, 45):
            import math
            ax = c + int((c - m * 2) * math.cos(math.radians(angle)))
            ay = c + int((c - m * 2) * math.sin(math.radians(angle)))
            draw.line([(c, c), (ax, ay)], fill=(r, g, b, 100), width=lw)
    elif any(k in name for k in ("dollar", "money", "currency", "calc")):
        # $ sign
        draw.text((m, m), "$", fill=(r, g, b, 220))
        draw.ellipse([(m, m * 2), (size - m, size - m * 2)], outline=(r, g, b, 220), width=lw)
    elif any(k in name for k in ("clock", "time", "hour", "timer")):
        # Clock
        draw.ellipse([(m, m), (size - m, size - m)], outline=(r, g, b, 220), width=lw)
        draw.line([(c, c), (c, m * 2)], fill=(r, g, b, 220), width=lw * 2)
        draw.line([(c, c), (c + m * 2, c)], fill=(r, g, b, 220), width=lw)
    elif any(k in name for k in ("star", "favorite")):
        # Star
        import math
        pts = []
        for i in range(5):
            a = math.radians(i * 72 - 90)
            pts.append((c + int((c - m) * math.cos(a)), c + int((c - m) * math.sin(a))))
            a2 = math.radians(i * 72 - 90 + 36)
            pts.append((c + int((c - m * 3) * math.cos(a2)), c + int((c - m * 3) * math.sin(a2))))
        draw.polygon(pts, fill=(r, g, b, 220))
    else:
        # Generic: filled circle with inner dot
        draw.ellipse([(m, m), (size - m, size - m)], outline=(r, g, b, 220), width=lw * 2)
        draw.ellipse([(c - m, c - m), (c + m, c + m)], fill=(r, g, b, 220))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ── Public functions used by pptx_visual_tools ────────────────────────────────

def get_icon_png_path(icon_name: str, size: int = 64, color_hex: str = "#1D4ED8") -> str | None:
    """Resolve icon name → local PNG temp file path.

    Returns a temp file path suitable for pptx slide.shapes.add_picture(),
    or None if the icon is unknown.
    """
    index = _load_index()
    filename = index.get(icon_name)
    if not filename:
        # Try fuzzy: any key that contains icon_name substring
        for k, v in index.items():
            if icon_name in k or k in icon_name:
                filename = v
                break
    if not filename:
        return None

    svg_path = _ICONS_DIR / filename
    if not svg_path.exists():
        return None

    png_bytes = _svg_to_png_bytes(svg_path, size, color_hex)
    if not png_bytes:
        return None

    fd, tmp_path = tempfile.mkstemp(suffix=".png")
    os.close(fd)
    with open(tmp_path, "wb") as f:
        f.write(png_bytes)
    return tmp_path


def get_placeholder_image_path() -> str:
    """Write the built-in base64 placeholder JPEG to a temp file and return its path."""
    raw = base64.b64decode(_PLACEHOLDER_B64)
    fd, tmp_path = tempfile.mkstemp(suffix=".jpg")
    os.close(fd)
    with open(tmp_path, "wb") as f:
        f.write(raw)
    return tmp_path


# ── Exposed tools ──────────────────────────────────────────────────────────────

@tool(description="List all locally bundled icons available for use in slides, grouped by category.")
def pptx_list_icons() -> ToolResult:
    """Return the full icon catalogue from the bundled assets/icons/ directory."""
    index = _load_index()
    if not index:
        return ToolResult.error(f"Icon index not found at {_INDEX_FILE}")

    categories = {
        "策略 & 目标": ["target", "crosshair", "flag", "trophy", "medal", "crown"],
        "增长 & 数据": ["chart_bar", "chart_line", "chart_pie", "chart_donut", "trend_up", "arrow_up_right"],
        "洞察 & 创新": ["lightbulb", "brain", "sparkle", "magic_wand", "eye", "search"],
        "执行 & 流程": ["check", "checks", "clock", "timer", "hourglass", "rocket", "lightning"],
        "人 & 组织":   ["users", "user", "user_circle", "buildings", "handshake", "megaphone"],
        "技术 & 数字": ["cpu", "database", "cloud", "code", "gear", "wrench"],
        "财务 & 商业": ["dollar", "trend_down", "calculator", "file_text", "clipboard", "receipt"],
        "连接 & 全球": ["globe", "map_pin", "network", "link", "wifi", "bell"],
        "安全 & 信任": ["shield", "lock", "key", "seal", "star", "heart"],
        "方向 & 导航": ["arrow_right", "expand", "compass", "path", "funnel", "stack"],
    }

    out = {}
    for cat, names in categories.items():
        available = [n for n in names if n in index]
        out[cat] = available

    return ToolResult.ok(json.dumps({
        "total": len(index),
        "categories": out,
        "usage_tip": "Pass icon name to layout content JSON, e.g. {'icon': 'target'} in three_cards layout",
    }, ensure_ascii=False))


@tool(
    description=(
        "Embed a locally bundled icon SVG into a PPTX slide at a specified position. "
        "icon_name must be one returned by pptx_list_icons(). "
        "size is in points (default 48). color is a hex string (default brand blue #1D4ED8)."
    )
)
def pptx_embed_icon(
    path: str,
    slide_index: int,
    icon_name: str,
    x: float,
    y: float,
    size: float = 48.0,
    color: str = "#1D4ED8",
) -> ToolResult:
    """Place a vector icon on an existing slide."""
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Pt  # type: ignore
    except ImportError:
        return ToolResult.error("python-pptx is not installed. Run: pip install python-pptx")

    p = Path(path).expanduser()
    if not p.exists():
        return ToolResult.error(f"File not found: {path}")

    prs = Presentation(str(p))
    if slide_index < 0 or slide_index >= len(prs.slides):
        return ToolResult.error(f"slide_index {slide_index} out of range (0–{len(prs.slides)-1})")

    slide = prs.slides[slide_index]
    icon_path = get_icon_png_path(icon_name, size=int(size * 2), color_hex=color)
    if not icon_path:
        return ToolResult.error(f"Icon '{icon_name}' not found. Use pptx_list_icons() to see available icons.")

    try:
        slide.shapes.add_picture(icon_path, Pt(x), Pt(y), Pt(size), Pt(size))
        prs.save(str(p))
    finally:
        try:
            os.unlink(icon_path)
        except Exception:
            pass

    return ToolResult.ok(json.dumps({
        "slide": slide_index,
        "icon": icon_name,
        "x": x, "y": y, "size": size,
    }, ensure_ascii=False))
