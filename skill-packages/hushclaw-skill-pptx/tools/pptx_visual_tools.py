"""Visual layout tools for high-quality PPT generation.

Implements a schema-driven rendering pipeline inspired by professional AI presentation
generators. Each layout function accepts a structured JSON payload and renders a
pixel-perfect slide using python-pptx primitives.

Layout catalogue (8 types):
  hero            – Full-bleed cover / section opener with bold headline
  image_split     – Left-image / right-text (or mirrored) split layout
  three_cards     – Three equal icon+title+body cards in a row
  big_stat        – 1–3 giant KPI numbers with labels
  quote           – Full-width pull-quote with accent bar and attribution
  two_column      – Two balanced text/bullet columns
  timeline        – Horizontal 3–5 node milestone bar
  agenda          – Numbered list / table-of-contents layout

Icons and placeholder images are sourced from local assets (assets/icons/) —
no external resource downloads at render time. pptx_fetch_image() calls the
Pexels API when online; falls back to the bundled placeholder JPEG when offline.
"""
from __future__ import annotations

import json
import urllib.request
import urllib.parse
from pathlib import Path

from hushclaw.tools.base import ToolResult, tool

# Import local icon/placeholder helpers (same skill package)
try:
    from tools.pptx_icon_tools import get_icon_png_path, get_placeholder_image_path  # type: ignore
except ImportError:
    try:
        from pptx_icon_tools import get_icon_png_path, get_placeholder_image_path  # type: ignore
    except ImportError:
        def get_icon_png_path(*a, **kw):  # type: ignore
            return None
        def get_placeholder_image_path():  # type: ignore
            return None

# ──────────────────────────────────────────────────────────────────────────────
# Internal colour palette (consulting_clean default; overridable per slide)
# ──────────────────────────────────────────────────────────────────────────────
_PALETTE = {
    "navy":       (15,  23,  42),   # #0F172A
    "blue":       (29,  78, 216),   # #1D4ED8
    "blue_light": (59, 130, 246),   # #3B82F6
    "slate":      (100, 116, 139),  # #64748B
    "mist":       (226, 232, 240),  # #E2E8F0
    "white":      (255, 255, 255),
    "black":      (0,   0,   0),
    "accent":     (99, 102, 241),   # indigo-500
    "green":      (16, 185, 129),   # emerald-500
    "amber":      (245, 158, 11),   # amber-500
    "red":        (239, 68,  68),   # red-500
}

# ──────────────────────────────────────────────────────────────────────────────
# Slide canvas constants  (1280 × 720 pt)
# ──────────────────────────────────────────────────────────────────────────────
_W = 1280  # slide width  in points
_H = 720   # slide height in points


def _rgb(name_or_hex: str | tuple):
    """Return an RGBColor from a palette name, hex string, or (r,g,b) tuple."""
    from pptx.dml.color import RGBColor  # type: ignore
    if isinstance(name_or_hex, tuple):
        return RGBColor(*name_or_hex)
    if isinstance(name_or_hex, str) and name_or_hex in _PALETTE:
        return RGBColor(*_PALETTE[name_or_hex])
    clean = name_or_hex.lstrip("#")
    return RGBColor(int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16))


def _pt(n):
    from pptx.util import Pt  # type: ignore
    return Pt(n)


def _in(n):
    from pptx.util import Inches  # type: ignore
    return Inches(n)


def _px(n):
    """Convert pixel value to EMU (1 pt = 12700 EMU; 1 px ≈ 1 pt at 96dpi)."""
    from pptx.util import Pt  # type: ignore
    return Pt(n)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill, line_color=None, line_w=0, radius=0):
    """Add a (optionally rounded) rectangle and return the shape."""
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, _px(x), _px(y), _px(w), _px(h))
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(fill)
    if line_color and line_w:
        s.line.color.rgb = _rgb(line_color)
        s.line.width = _pt(line_w)
    else:
        s.line.fill.background()
    if radius:
        try:
            s.adjustments[0] = radius / min(w, h)
        except Exception:
            pass
    return s


def _textbox(slide, x, y, w, h, text, size, color, bold=False, italic=False,
             align="left", font="Calibri", wrap=True, line_spacing=None):
    """Add a text box and return the shape."""
    from pptx.enum.text import PP_ALIGN  # type: ignore
    box = slide.shapes.add_textbox(_px(x), _px(y), _px(w), _px(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    if line_spacing:
        p.line_spacing = line_spacing
    _align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    p.alignment = _align_map.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = _pt(size)
    run.font.color.rgb = _rgb(color)
    run.font.bold = bold
    run.font.italic = italic
    return box


def _multiline_textbox(slide, x, y, w, h, lines: list[dict], wrap=True):
    """Add a text box with multiple paragraph/run specs.
    Each dict: {text, size, color, bold=False, italic=False, align='left', space_before=0}
    """
    from pptx.enum.text import PP_ALIGN  # type: ignore
    box = slide.shapes.add_textbox(_px(x), _px(y), _px(w), _px(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    _align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if spec.get("space_before"):
            p.space_before = _pt(spec["space_before"])
        p.alignment = _align_map.get(spec.get("align", "left"), PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = spec.get("text", "")
        run.font.name = spec.get("font", "Calibri")
        run.font.size = _pt(spec.get("size", 16))
        run.font.color.rgb = _rgb(spec.get("color", "white"))
        run.font.bold = spec.get("bold", False)
        run.font.italic = spec.get("italic", False)
    return box


def _accent_bar(slide, x, y, h, color="blue"):
    _rect(slide, x, y, 5, h, color)


def _is_emoji(s: str) -> bool:
    """Return True if s looks like a single emoji character (codepoint ≥ U+1F000 or common symbols)."""
    if not s:
        return False
    # Emoji are typically non-ASCII, single code-unit or surrogate pair
    return len(s) <= 4 and any(ord(c) > 0x2500 for c in s)


def _parse_json_arg(raw: str) -> tuple[dict | None, str]:
    try:
        data = json.loads(raw)
        if not isinstance(data, dict):
            return None, "content_json must be a JSON object"
        return data, ""
    except json.JSONDecodeError as e:
        return None, f"Invalid JSON: {e}"


def _try_import_pptx():
    try:
        from pptx import Presentation  # type: ignore
        return Presentation, None
    except ImportError:
        return None, "python-pptx is not installed. Run: pip install python-pptx"


# ──────────────────────────────────────────────────────────────────────────────
# Tool: fetch image from Pexels (free, no auth for basic URL format)
# ──────────────────────────────────────────────────────────────────────────────

@tool(
    description=(
        "Fetch a relevant high-quality image URL from Pexels for use in a slide. "
        "Returns a direct image URL you can pass to image_url in layout tools. "
        "Falls back to the bundled local placeholder image (no network) when Pexels is unavailable."
    )
)
def pptx_fetch_image(
    query: str,
    orientation: str = "landscape",
) -> ToolResult:
    """Search Pexels for a photo matching the query and return a direct image URL.

    When offline or the API key is invalid, returns 'local://placeholder' which
    the rendering pipeline will resolve to the bundled gray-gradient placeholder.
    """
    query = str(query or "business").strip()
    try:
        encoded = urllib.parse.quote(query)
        url = (
            "https://api.pexels.com/v1/search"
            f"?query={encoded}&orientation={orientation}&per_page=5&page=1"
        )
        req = urllib.request.Request(
            url,
            headers={"Authorization": "563492ad6f9170000100000149a6f8b5bdd94fd6b7568fc4e11d6d88"},
        )
        with urllib.request.urlopen(req, timeout=8) as resp:
            data = json.loads(resp.read())
            photos = data.get("photos", [])
            if photos:
                img_url = photos[0]["src"]["large2x"]
                return ToolResult.ok(json.dumps({"image_url": img_url, "source": "pexels", "query": query}, ensure_ascii=False))
    except Exception:
        pass

    # Offline / API failure — use the bundled local placeholder (no external URL needed)
    return ToolResult.ok(json.dumps({
        "image_url": "local://placeholder",
        "source": "local_placeholder",
        "query": query,
    }, ensure_ascii=False))


# ──────────────────────────────────────────────────────────────────────────────
# Tool: list available layouts with their content schemas
# ──────────────────────────────────────────────────────────────────────────────

_LAYOUT_CATALOGUE = {
    "hero": {
        "description": "Full-bleed cover / section opener. Bold headline + subtitle + optional tag line.",
        "best_for": "Opening slide, chapter divider, key statement",
        "schema": {
            "title": "string (max 30 chars) — punchy bold headline",
            "subtitle": "string (max 80 chars) — supporting sentence",
            "tag": "string (max 20 chars, optional) — e.g. 'Q2 2026' or 'CONFIDENTIAL'",
            "accent_color": "hex or palette name (optional, default: blue)",
            "bg_color": "hex or palette name (optional, default: navy)",
        },
    },
    "image_split": {
        "description": "Left image, right text with headline + bullet points. Mirror with side='right'.",
        "best_for": "Presenting a concept with visual evidence; product/service features",
        "schema": {
            "title": "string (max 35 chars)",
            "bullets": "array of strings, max 4 items, each max 60 chars",
            "image_url": "URL string — use pptx_fetch_image() to obtain",
            "side": "'left' or 'right' (optional, default: left) — which side the image is on",
            "so_what": "string (max 80 chars, optional) — action/implication callout",
        },
    },
    "three_cards": {
        "description": "Three equal icon+title+body cards in a row below a section headline.",
        "best_for": "Comparing 3 options, 3 pillars, 3 findings",
        "schema": {
            "title": "string (max 35 chars)",
            "cards": "array of exactly 3 objects: {icon: emoji or symbol, heading: max 25 chars, body: max 80 chars}",
            "accent_color": "hex or palette name (optional)",
        },
    },
    "big_stat": {
        "description": "1–3 giant KPI numbers with labels and optional delta indicators.",
        "best_for": "Quantitative highlights, ROI slides, benchmark comparisons",
        "schema": {
            "title": "string (max 40 chars)",
            "stats": "array of 1–3 objects: {value: '22%', label: max 30 chars, delta: '+3pp vs LY' optional, color: optional}",
            "footnote": "string (max 80 chars, optional) — data source / period",
        },
    },
    "quote": {
        "description": "Full-width pull-quote with large text, accent bar, and attribution.",
        "best_for": "Customer voice, executive quote, key insight statement",
        "schema": {
            "quote": "string (max 160 chars) — the quote text",
            "author": "string (max 50 chars) — name + title",
            "accent_color": "hex or palette name (optional)",
        },
    },
    "two_column": {
        "description": "Two balanced columns, each with a headline + bullet list.",
        "best_for": "Before/after, pros/cons, two workstreams, compare two options",
        "schema": {
            "title": "string (max 40 chars)",
            "left_heading": "string (max 25 chars)",
            "left_bullets": "array of strings, max 5 items, each max 60 chars",
            "right_heading": "string (max 25 chars)",
            "right_bullets": "array of strings, max 5 items, each max 60 chars",
            "left_color": "hex or palette name (optional, default: blue)",
            "right_color": "hex or palette name (optional, default: accent)",
        },
    },
    "timeline": {
        "description": "Horizontal 3–5 node milestone bar with dates and descriptions.",
        "best_for": "Project roadmap, implementation phases, historical journey",
        "schema": {
            "title": "string (max 40 chars)",
            "nodes": "array of 3–5 objects: {date: max 12 chars, label: max 20 chars, body: max 60 chars}",
            "accent_color": "hex or palette name (optional)",
        },
    },
    "agenda": {
        "description": "Numbered agenda / table-of-contents layout with section headings.",
        "best_for": "Opening agenda, section recap, structured outline",
        "schema": {
            "title": "string (max 30 chars, default: 'Agenda')",
            "items": "array of 3–7 objects: {number: '01', heading: max 30 chars, desc: max 60 chars optional}",
            "accent_color": "hex or palette name (optional)",
        },
    },
}


@tool(description="List all available visual slide layouts with their content schemas and use cases.")
def pptx_list_visual_layouts() -> ToolResult:
    """Return catalogue of all 8 layouts with schemas for schema-driven rendering."""
    items = []
    for name, meta in _LAYOUT_CATALOGUE.items():
        items.append({
            "layout_id": name,
            "description": meta["description"],
            "best_for": meta["best_for"],
            "schema": meta["schema"],
        })
    return ToolResult.ok(json.dumps({"layouts": items, "total": len(items)}, ensure_ascii=False))


# ──────────────────────────────────────────────────────────────────────────────
# Tool: unified schema-driven renderer
# ──────────────────────────────────────────────────────────────────────────────

@tool(
    description=(
        "Add a professionally designed slide using a named layout and structured JSON content. "
        "layout_id must be one of: hero, image_split, three_cards, big_stat, quote, two_column, timeline, agenda. "
        "Use pptx_list_visual_layouts() to get the exact schema for each layout. "
        "This is the primary rendering tool for high-quality PPT generation."
    )
)
def pptx_add_visual_slide(
    path: str,
    layout_id: str,
    content_json: str,
) -> ToolResult:
    """Render a high-quality slide from a layout id and content JSON string."""
    Presentation, err = _try_import_pptx()
    if err:
        return ToolResult.error(err)

    p = Path(path).expanduser()
    if not p.exists():
        return ToolResult.error(f"File not found: {path}")

    data, err = _parse_json_arg(content_json)
    if err:
        return ToolResult.error(err)

    layout_id = str(layout_id).strip().lower()
    if layout_id not in _LAYOUT_CATALOGUE:
        return ToolResult.error(f"Unknown layout '{layout_id}'. Available: {list(_LAYOUT_CATALOGUE)}")

    prs = Presentation(str(p))
    prs.slide_width = _px(_W)
    prs.slide_height = _px(_H)

    dispatch = {
        "hero": _render_hero,
        "image_split": _render_image_split,
        "three_cards": _render_three_cards,
        "big_stat": _render_big_stat,
        "quote": _render_quote,
        "two_column": _render_two_column,
        "timeline": _render_timeline,
        "agenda": _render_agenda,
    }
    try:
        slide_index = dispatch[layout_id](prs, data)
    except Exception as exc:
        import traceback
        return ToolResult.error(f"Render error in '{layout_id}': {exc}\n{traceback.format_exc()}")

    prs.save(str(p))
    return ToolResult.ok(json.dumps({
        "added_slide": slide_index,
        "layout": layout_id,
        "path": str(p),
    }, ensure_ascii=False))


# ──────────────────────────────────────────────────────────────────────────────
# Layout renderers
# ──────────────────────────────────────────────────────────────────────────────

def _render_hero(prs, d: dict) -> int:
    slide = _blank_slide(prs)
    bg_color = d.get("bg_color", "navy")
    accent_color = d.get("accent_color", "blue")

    # Full background
    _rect(slide, 0, 0, _W, _H, bg_color)

    # Decorative gradient bar left edge
    _rect(slide, 0, 0, 8, _H, accent_color)

    # Large accent circle (decorative, top-right)
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore
    circ = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        _px(_W - 320), _px(-120), _px(360), _px(360),
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = _rgb(accent_color)
    circ.line.fill.background()
    # Low opacity via XML alpha
    try:
        from lxml import etree  # type: ignore
        sp_pr = circ._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if sp_pr is not None:
            clr = sp_pr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if clr is not None:
                alpha_el = etree.SubElement(clr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                alpha_el.set('val', '15000')  # ~15% opacity
    except Exception:
        pass

    # Optional tag  (top-left pill)
    tag = str(d.get("tag", "")).strip()
    if tag:
        _rect(slide, 55, 48, len(tag) * 11 + 40, 36, accent_color, radius=18)
        _textbox(slide, 70, 51, len(tag) * 11 + 20, 30, tag,
                 size=13, color="white", bold=True, align="left", font="Calibri")

    # Title
    title = str(d.get("title", "Headline")).strip()[:60]
    _textbox(slide, 55, 120, 900, 200, title,
             size=68, color="white", bold=True, font="Calibri", align="left")

    # Horizontal rule
    _rect(slide, 55, 320, 120, 4, accent_color)

    # Subtitle
    subtitle = str(d.get("subtitle", "")).strip()[:120]
    if subtitle:
        _textbox(slide, 55, 340, 900, 120, subtitle,
                 size=26, color=(148, 163, 184), bold=False, font="Calibri")

    return len(prs.slides) - 1


def _render_image_split(prs, d: dict) -> int:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore
    slide = _blank_slide(prs)

    side = str(d.get("side", "left")).strip().lower()
    img_x = 0 if side == "left" else 660
    txt_x = 660 if side == "left" else 30

    # Background
    _rect(slide, 0, 0, _W, _H, "white")
    # Image placeholder / loaded image
    image_url = str(d.get("image_url", "")).strip()
    img_panel = _rect(slide, img_x, 0, 620, _H, (226, 232, 240))
    if image_url:
        img_path = _download_image_to_temp(image_url)
        if img_path:
            try:
                slide.shapes.add_picture(img_path, _px(img_x), _px(0), _px(620), _px(_H))
            except Exception:
                pass  # keep the grey placeholder

    # Accent strip between panels
    _rect(slide, img_x + (620 if side == "left" else -8), 0, 8, _H, "blue")

    # Title
    title = str(d.get("title", "Title")).strip()[:50]
    _textbox(slide, txt_x + 25, 80, 560, 100, title,
             size=36, color="navy", bold=True, font="Calibri")

    # Bullets
    bullets = d.get("bullets") or []
    if isinstance(bullets, list):
        by = 200
        for i, bul in enumerate(bullets[:4]):
            bul_text = str(bul).strip()[:80]
            # bullet dot
            dot = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                _px(txt_x + 25), _px(by + 6), _px(10), _px(10),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = _rgb("blue")
            dot.line.fill.background()
            _textbox(slide, txt_x + 48, by, 520, 60, bul_text,
                     size=18, color="navy", font="Calibri")
            by += 70

    # So-what box
    so_what = str(d.get("so_what", "")).strip()[:100]
    if so_what:
        _rect(slide, txt_x + 20, 570, 580, 100, (239, 246, 255), "blue", 1, radius=4)
        _multiline_textbox(slide, txt_x + 36, 580, 548, 80, [
            {"text": "Key takeaway", "size": 12, "color": "blue", "bold": True},
            {"text": so_what, "size": 15, "color": "navy", "space_before": 4},
        ])

    return len(prs.slides) - 1


def _render_three_cards(prs, d: dict) -> int:
    slide = _blank_slide(prs)
    accent = d.get("accent_color", "blue")

    # Background
    _rect(slide, 0, 0, _W, _H, "white")

    # Top accent bar
    _rect(slide, 0, 0, _W, 6, accent)

    # Section title
    title = str(d.get("title", "")).strip()[:50]
    _textbox(slide, 50, 40, 1180, 80, title,
             size=34, color="navy", bold=True, font="Calibri")

    # Divider
    _rect(slide, 50, 120, 1180, 1, (226, 232, 240))

    # Three cards
    cards = d.get("cards") or []
    if not isinstance(cards, list):
        cards = []
    cards = cards[:3]
    while len(cards) < 3:
        cards.append({"icon": "◆", "heading": "Point", "body": ""})

    card_w = 360
    card_h = 430
    gap = ((_W - 100) - card_w * 3) // 2
    cx = 50

    for card in cards:
        icon = str(card.get("icon", "◆"))
        heading = str(card.get("heading", "")).strip()[:30]
        body = str(card.get("body", "")).strip()[:120]

        # Card shadow effect (thin border)
        _rect(slide, cx + 3, 153, card_w, card_h, (203, 213, 225), radius=6)
        # Card body
        card_shape = _rect(slide, cx, 150, card_w, card_h, "white",
                           (214, 224, 255), 1, radius=6)

        # Accent top
        _rect(slide, cx, 150, card_w, 6, accent, radius=3)

        # Icon circle background
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore
        ic_circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            _px(cx + 28), _px(198), _px(68), _px(68),
        )
        ic_circle.fill.solid()
        ic_circle.fill.fore_color.rgb = _rgb(accent)
        ic_circle.line.fill.background()

        # Try to embed a local PNG icon; fall back to emoji text
        accent_hex = "#" + "".join(f"{v:02X}" for v in _PALETTE.get(accent, _PALETTE["blue"])) if accent in _PALETTE else str(accent)
        icon_png = get_icon_png_path(icon, size=96, color_hex="#FFFFFF") if icon and len(icon) > 1 and not _is_emoji(icon) else None
        if icon_png:
            try:
                slide.shapes.add_picture(icon_png, _px(cx + 36), _px(206), _px(52), _px(52))
            except Exception:
                _textbox(slide, cx + 28, 202, 68, 60, icon,
                         size=28, color="white", bold=True, align="center", font="Segoe UI Emoji")
            finally:
                import os as _os
                try:
                    _os.unlink(icon_png)
                except Exception:
                    pass
        else:
            _textbox(slide, cx + 28, 202, 68, 60, icon,
                     size=28, color="white", bold=True, align="center", font="Segoe UI Emoji")

        # Heading
        _textbox(slide, cx + 20, 285, card_w - 40, 60, heading,
                 size=22, color="navy", bold=True, font="Calibri")

        # Body
        _textbox(slide, cx + 20, 350, card_w - 40, 200, body,
                 size=16, color=(71, 85, 105), font="Calibri", wrap=True)

        cx += card_w + gap

    return len(prs.slides) - 1


def _render_big_stat(prs, d: dict) -> int:
    slide = _blank_slide(prs)

    _rect(slide, 0, 0, _W, _H, "navy")

    # Title
    title = str(d.get("title", "")).strip()[:60]
    _textbox(slide, 50, 50, 1180, 80, title,
             size=32, color=(148, 163, 184), bold=False, font="Calibri", align="center")

    # Stats
    stats = d.get("stats") or []
    if not isinstance(stats, list):
        stats = []
    stats = stats[:3]
    if not stats:
        stats = [{"value": "—", "label": "KPI"}]

    n = len(stats)
    col_w = _W // n
    stat_colors = ["blue_light", "green", "amber"]

    for i, stat in enumerate(stats):
        value = str(stat.get("value", "—")).strip()[:12]
        label = str(stat.get("label", "")).strip()[:35]
        delta = str(stat.get("delta", "")).strip()[:25]
        sc = stat.get("color", stat_colors[i % len(stat_colors)])

        cx = i * col_w
        # Vertical divider (except first)
        if i > 0:
            _rect(slide, cx, 160, 1, 360, (30, 41, 59))

        # Value
        _textbox(slide, cx + 20, 150, col_w - 40, 220, value,
                 size=110, color=sc, bold=True, font="Calibri", align="center")

        # Label
        _textbox(slide, cx + 20, 370, col_w - 40, 60, label,
                 size=20, color=(148, 163, 184), font="Calibri", align="center")

        # Delta
        if delta:
            delta_color = "green" if delta.startswith("+") else (
                "red" if delta.startswith("-") else (148, 163, 184)
            )
            _textbox(slide, cx + 20, 430, col_w - 40, 40, delta,
                     size=16, color=delta_color, bold=True, align="center", font="Calibri")

    # Footnote
    footnote = str(d.get("footnote", "")).strip()[:100]
    if footnote:
        _textbox(slide, 50, 640, 1180, 40, footnote,
                 size=13, color=(71, 85, 105), font="Calibri", align="center")

    return len(prs.slides) - 1


def _render_quote(prs, d: dict) -> int:
    slide = _blank_slide(prs)
    accent = d.get("accent_color", "blue")

    _rect(slide, 0, 0, _W, _H, "navy")

    # Large decorative quotation mark
    _textbox(slide, 60, 60, 200, 180, "\u201c",
             size=200, color=(29, 78, 216), bold=True, font="Georgia", align="left")

    # Accent bar
    _rect(slide, 60, 250, 6, 220, accent)

    # Quote text
    quote = str(d.get("quote", "")).strip()[:200]
    _textbox(slide, 90, 240, 1100, 280, quote,
             size=34, color="white", bold=False, font="Georgia", wrap=True)

    # Attribution
    author = str(d.get("author", "")).strip()[:60]
    _textbox(slide, 90, 550, 1100, 60, f"— {author}",
             size=18, color=(100, 116, 139), bold=False, font="Calibri")

    # Bottom rule
    _rect(slide, 90, 620, 300, 2, accent)

    return len(prs.slides) - 1


def _render_two_column(prs, d: dict) -> int:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore
    slide = _blank_slide(prs)

    _rect(slide, 0, 0, _W, _H, "white")
    _rect(slide, 0, 0, _W, 6, "blue")

    # Title
    title = str(d.get("title", "")).strip()[:60]
    _textbox(slide, 50, 30, 1180, 70, title,
             size=32, color="navy", bold=True, font="Calibri")

    _rect(slide, 50, 100, 1180, 1, (226, 232, 240))

    # Columns setup
    left_color = d.get("left_color", "blue")
    right_color = d.get("right_color", "accent")
    col_w = 545

    for col_idx, (heading_key, bullets_key, col_color, cx) in enumerate([
        ("left_heading", "left_bullets", left_color, 50),
        ("right_heading", "right_bullets", right_color, 665),
    ]):
        heading = str(d.get(heading_key, "") or "").strip()[:30]
        bullets = d.get(bullets_key) or []
        if not isinstance(bullets, list):
            bullets = []

        # Column header bar
        _rect(slide, cx, 120, col_w, 48, col_color, radius=4)
        _textbox(slide, cx + 15, 126, col_w - 30, 36, heading,
                 size=20, color="white", bold=True, font="Calibri")

        # Bullets
        by = 190
        for bul in bullets[:5]:
            bul_text = str(bul).strip()[:70]
            # bullet
            dot = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                _px(cx + 10), _px(by + 7), _px(8), _px(8),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = _rgb(col_color)
            dot.line.fill.background()
            _textbox(slide, cx + 30, by, col_w - 35, 58, bul_text,
                     size=17, color=(30, 41, 59), font="Calibri", wrap=True)
            by += 70

    # Center divider
    _rect(slide, 635, 120, 1, 550, (226, 232, 240))

    return len(prs.slides) - 1


def _render_timeline(prs, d: dict) -> int:
    slide = _blank_slide(prs)
    accent = d.get("accent_color", "blue")
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore

    _rect(slide, 0, 0, _W, _H, "white")
    _rect(slide, 0, 0, _W, 6, accent)

    # Title
    title = str(d.get("title", "")).strip()[:60]
    _textbox(slide, 50, 30, 1180, 70, title,
             size=32, color="navy", bold=True, font="Calibri")

    nodes = d.get("nodes") or []
    if not isinstance(nodes, list):
        nodes = []
    nodes = nodes[:5]
    n = max(len(nodes), 1)

    # Horizontal line
    line_y = 350
    _rect(slide, 80, line_y - 2, _W - 160, 4, (203, 213, 225))

    # Node spacing
    spacing = (_W - 160) / (n - 1) if n > 1 else 0
    start_x = 80

    for i, node in enumerate(nodes):
        nx = int(start_x + i * spacing)
        date = str(node.get("date", "")).strip()[:15]
        label = str(node.get("label", "")).strip()[:25]
        body = str(node.get("body", "")).strip()[:80]

        # Connector dot
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            _px(nx - 18), _px(line_y - 18), _px(36), _px(36),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(accent)
        dot.line.fill.background()
        # Inner white dot
        inner = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            _px(nx - 8), _px(line_y - 8), _px(16), _px(16),
        )
        inner.fill.solid()
        inner.fill.fore_color.rgb = _rgb("white")
        inner.line.fill.background()

        # Vertical stem (alternating above/below)
        above = (i % 2 == 0)
        stem_top = line_y - 80 if above else line_y + 36
        _rect(slide, nx - 1, min(stem_top, line_y - 18), 2, 62, (203, 213, 225))

        if above:
            # Date above line
            _textbox(slide, nx - 80, line_y - 220, 160, 40, date,
                     size=13, color=accent, bold=True, align="center", font="Calibri")
            # Label
            _textbox(slide, nx - 100, line_y - 175, 200, 45, label,
                     size=17, color="navy", bold=True, align="center", font="Calibri")
            # Body
            _textbox(slide, nx - 120, line_y - 130, 240, 65, body,
                     size=13, color=(71, 85, 105), align="center", font="Calibri", wrap=True)
        else:
            # Date below line
            _textbox(slide, nx - 80, line_y + 110, 160, 40, date,
                     size=13, color=accent, bold=True, align="center", font="Calibri")
            _textbox(slide, nx - 100, line_y + 150, 200, 45, label,
                     size=17, color="navy", bold=True, align="center", font="Calibri")
            _textbox(slide, nx - 120, line_y + 195, 240, 65, body,
                     size=13, color=(71, 85, 105), align="center", font="Calibri", wrap=True)

    return len(prs.slides) - 1


def _render_agenda(prs, d: dict) -> int:
    slide = _blank_slide(prs)
    accent = d.get("accent_color", "blue")

    _rect(slide, 0, 0, _W, _H, "white")

    # Left accent panel
    _rect(slide, 0, 0, 320, _H, "navy")

    # Title on left panel
    title = str(d.get("title", "Agenda")).strip()[:25]
    _textbox(slide, 30, 200, 260, 300, title,
             size=52, color="white", bold=True, font="Calibri", wrap=True)
    _rect(slide, 50, 500, 80, 5, accent)

    # Items
    items = d.get("items") or []
    if not isinstance(items, list):
        items = []
    items = items[:7]

    row_h = min(80, 560 // max(len(items), 1))
    iy = ((_H - row_h * len(items)) // 2)

    for j, item in enumerate(items):
        num = str(item.get("number", f"{j+1:02d}")).strip()
        heading = str(item.get("heading", "")).strip()[:35]
        desc = str(item.get("desc", "")).strip()[:60]

        row_y = iy + j * row_h

        # Number badge
        _rect(slide, 355, row_y, 48, 48, accent, radius=4)
        _textbox(slide, 355, row_y + 4, 48, 40, num,
                 size=18, color="white", bold=True, align="center", font="Calibri")

        # Heading + desc
        _textbox(slide, 420, row_y, 820, 36, heading,
                 size=22, color="navy", bold=True, font="Calibri")
        if desc:
            _textbox(slide, 420, row_y + 32, 820, 36, desc,
                     size=15, color=(100, 116, 139), font="Calibri")

        # Row divider
        if j < len(items) - 1:
            _rect(slide, 355, row_y + row_h - 4, 870, 1, (226, 232, 240))

    return len(prs.slides) - 1


# ──────────────────────────────────────────────────────────────────────────────
# Helper: download image to a temp file for embedding
# ──────────────────────────────────────────────────────────────────────────────

def _download_image_to_temp(url: str) -> str | None:
    """Download an image URL to a temp file and return the path.

    Handles three cases:
    - 'local://placeholder'  → write bundled base64 placeholder to temp file
    - 'file://...' or absolute path → read from disk
    - HTTP(S) URL → download with timeout; falls back to local placeholder on error
    """
    import tempfile
    import os

    # Case 1: explicit local placeholder sentinel
    if url == "local://placeholder":
        return get_placeholder_image_path()

    # Case 2: local file path / file:// URI
    if url.startswith("file://"):
        return url[7:]
    if url and not url.startswith("http"):
        return url if os.path.exists(url) else get_placeholder_image_path()

    # Case 3: remote URL
    try:
        suffix = ".png" if ".png" in url.lower() else ".jpg"
        fd, tmp_path = tempfile.mkstemp(suffix=suffix)
        os.close(fd)
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=10) as resp:
            with open(tmp_path, "wb") as f:
                f.write(resp.read())
        return tmp_path
    except Exception:
        # Network failure → fall back to bundled placeholder
        return get_placeholder_image_path()


# ──────────────────────────────────────────────────────────────────────────────
# Convenience tool: build an entire deck from a slides spec list in one call
# ──────────────────────────────────────────────────────────────────────────────

@tool(
    description=(
        "Build a complete PPT deck from a JSON array of slide specs in one call. "
        "Each spec must have 'layout_id' and 'content' fields. "
        "Creates the file first (or uses existing), then renders each slide in order. "
        "Returns the final slide count and path."
    )
)
def pptx_build_visual_deck(
    path: str,
    slides_json: str,
    overwrite: bool = False,
) -> ToolResult:
    """One-shot deck builder: create file + render all slides from a spec list.

    slides_json format:
    [
      {"layout_id": "hero",    "content": {...}},
      {"layout_id": "image_split", "content": {...}},
      ...
    ]
    """
    Presentation, err = _try_import_pptx()
    if err:
        return ToolResult.error(err)

    try:
        slides_spec = json.loads(slides_json)
        if not isinstance(slides_spec, list):
            return ToolResult.error("slides_json must be a JSON array")
    except json.JSONDecodeError as e:
        return ToolResult.error(f"Invalid JSON in slides_json: {e}")

    p = Path(path).expanduser()
    p.parent.mkdir(parents=True, exist_ok=True)

    if not p.exists() or overwrite:
        prs = Presentation()
        prs.slide_width = _px(_W)
        prs.slide_height = _px(_H)
        prs.save(str(p))

    errors = []
    for i, spec in enumerate(slides_spec):
        layout_id = str(spec.get("layout_id", "")).strip().lower()
        content = spec.get("content", {})
        if not layout_id:
            errors.append(f"Slide {i}: missing layout_id")
            continue
        if layout_id not in _LAYOUT_CATALOGUE:
            errors.append(f"Slide {i}: unknown layout '{layout_id}'")
            continue
        content_str = json.dumps(content, ensure_ascii=False)
        result = pptx_add_visual_slide(path=str(p), layout_id=layout_id, content_json=content_str)
        if result.is_error:
            errors.append(f"Slide {i} ({layout_id}): {result.content}")

    # Final count
    prs = Presentation(str(p))
    final_count = len(prs.slides)

    return ToolResult.ok(json.dumps({
        "path": str(p),
        "slide_count": final_count,
        "errors": errors,
        "success": len(errors) == 0,
    }, ensure_ascii=False))
